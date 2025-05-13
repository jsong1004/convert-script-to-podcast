import os
import io
from flask import Blueprint, render_template, request, current_app
import google.generativeai as genai
from pptx import Presentation
import PyPDF2

# Ensure GOOGLE_API_KEY is loaded. genai should be configured in app.py
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")

presentation_bp = Blueprint('presentation_bp', __name__, template_folder='../templates')

def extract_text_from_pptx(pptx_file_stream):
    """Extracts all text from a .pptx file stream."""
    try:
        prs = Presentation(pptx_file_stream)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        return "\n".join(text_runs)
    except Exception as e:
        current_app.logger.error(f"Error extracting text from PPTX: {e}")
        raise ValueError(f"Could not extract text from presentation: {e}")

def extract_text_from_pdf(pdf_file_stream):
    """Extracts all text from a PDF file stream."""
    try:
        pdf_reader = PyPDF2.PdfReader(pdf_file_stream)
        text = []
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            text.append(page.extract_text())
        return "\n".join(text)
    except Exception as e:
        current_app.logger.error(f"Error extracting text from PDF: {e}")
        raise ValueError(f"Could not extract text from PDF: {e}")

def generate_script_with_gemini(presentation_text, script_style, output_language='en'):
    """Generates script using Google Gemini API with chunking and overlapping content."""
    if not GOOGLE_API_KEY:
        current_app.logger.error("Google API Key is not configured for Gemini.")
        raise ValueError("Google API Key is not configured.")

    if not presentation_text or not isinstance(presentation_text, str):
        raise ValueError("Invalid input text provided.")

    # Use GOOGLE_MODEL from .env or fallback
    env_model_name = os.getenv("GOOGLE_MODEL")
    primary_model_name = env_model_name if env_model_name else "gemini-1.5-pro"
    fallback_model_name = "gemini-2.0-flash"

    model_name = primary_model_name
    try:
        model = genai.GenerativeModel(model_name)
        current_app.logger.info(f"Successfully initialized Gemini model: {model_name}")
    except Exception as e:
        current_app.logger.error(f"Error initializing Gemini model '{model_name}': {e}")
        if model_name != fallback_model_name:
            current_app.logger.info(f"Attempting to use fallback model: {fallback_model_name}")
            try:
                model_name = fallback_model_name
                model = genai.GenerativeModel(model_name)
                current_app.logger.info(f"Successfully initialized fallback Gemini model: {model_name}")
            except Exception as fallback_e:
                current_app.logger.error(f"Error initializing fallback Gemini model '{model_name}': {fallback_e}")
                raise ValueError(f"Could not initialize Gemini model. Error: {fallback_e}")
        else:
            raise ValueError(f"Could not initialize Gemini model '{model_name}'. Error: {e}")

    # Split text into chunks with overlap
    def split_text_into_chunks_with_overlap(text, chunk_size=6000, overlap_size=1000):
        if not text or not isinstance(text, str):
            raise ValueError("Invalid text input for chunking")
            
        words = text.split()
        chunks = []
        current_chunk = []
        current_size = 0
        overlap_words = []
        
        for word in words:
            word_size = len(word) + 1  # +1 for space
            
            # If adding this word would exceed chunk size
            if current_size + word_size > chunk_size:
                # Save current chunk
                chunks.append(' '.join(current_chunk))
                
                # Keep last few words for overlap
                overlap_words = current_chunk[-overlap_size//10:]  # Approximate number of words for overlap
                
                # Start new chunk with overlap
                current_chunk = overlap_words + [word]
                current_size = sum(len(w) + 1 for w in current_chunk)
            else:
                current_chunk.append(word)
                current_size += word_size
        
        # Add the last chunk if it's not empty
        if current_chunk:
            chunks.append(' '.join(current_chunk))
        
        return chunks

    try:
        chunks = split_text_into_chunks_with_overlap(presentation_text)
        current_app.logger.info(f"Split input into {len(chunks)} chunks with overlap")
    except Exception as e:
        current_app.logger.error(f"Error splitting text into chunks: {e}")
        raise ValueError(f"Failed to process input text: {str(e)}")

    language_instruction = {
        'en': 'Output the script in English.',
        'ko': 'Output the script in Korean. 모든 결과를 한국어로 출력하세요.'
    }.get(output_language, 'Output the script in English.')

    common_prompt_instructions = (
        f"You are an expert scriptwriter. Convert the following presentation slide content into a spoken script. "
        f"The script should be very detailed for each slide concept. "
        f"The tone should be passionate, knowledgeable, and educational, like someone speaking engagingly in front of many people. "
        f"Include easy-to-understand examples where appropriate to clarify concepts. "
        f"Do not include any titles, subtitles, or descriptive slide markers like 'Slide 1 of 5' or 'Next slide'. "
        f"Only output the text that will be spoken. Ensure smooth transitions between ideas if they were on different slides. "
        f"{language_instruction}"
    )

    if script_style == "podcast":
        style_specific_instructions = (
            "The script should be in a podcast format with multiple voices. "
            "Use 'HOST:' for the main narrator and overall flow. "
            "Use 'VOICE 1:', 'VOICE 2:', etc., for different perspectives, examples, or to break up longer segments of HOST narration. "
            "Ensure a conversational and engaging flow between speakers. The HOST should guide the conversation."
        )
    elif script_style == "speech":
        style_specific_instructions = (
            "The script should be for a single speaker. All text should be attributed to 'HOST:'. "
            "Ensure a continuous, engaging monologue suitable for a keynote presentation."
        )
    else:
        raise ValueError("Invalid script style selected.")

    generated_scripts = []
    
    for i, chunk in enumerate(chunks):
        try:
            # Add context about chunk position and overlap
            chunk_context = (
                f"This is part {i+1} of {len(chunks)} of the presentation content. "
                f"Some content may overlap with the previous or next part to maintain context. "
                f"Focus on generating a coherent script for this section while maintaining continuity."
            )
            
            chunk_prompt = (
                f"{common_prompt_instructions}\n\n"
                f"{style_specific_instructions}\n\n"
                f"{chunk_context}\n\n"
                f"PRESENTATION CONTENT:\n---\n{chunk}\n---\n\n"
                f"GENERATED SCRIPT:"
            )
            
            current_app.logger.info(f"Processing chunk {i+1}/{len(chunks)}. Length: {len(chunk_prompt)}")
            
            response = model.generate_content(
                chunk_prompt,
                generation_config={
                    "temperature": 0.7,
                    "top_p": 0.8,
                    "top_k": 40,
                    "max_output_tokens": 2048,
                }
            )
            
            if hasattr(response, 'prompt_feedback') and response.prompt_feedback.block_reason:
                current_app.logger.error(f"Gemini content generation blocked. Reason: {response.prompt_feedback.block_reason}")
                current_app.logger.error(f"Block details: {response.prompt_feedback.safety_ratings}")
                raise ValueError(f"Content generation blocked by safety filters. Reason: {response.prompt_feedback.block_reason}")

            if not response.parts:
                current_app.logger.warning(f"Gemini response has no parts. Full response: {response}")
                if hasattr(response, 'text') and response.text:
                    generated_scripts.append(response.text.strip())
                else:
                    raise ValueError("Received an empty response from Gemini.")
            else:
                generated_text = "".join(part.text for part in response.parts if hasattr(part, 'text'))
                if not generated_text.strip():
                    current_app.logger.warning(f"Gemini generated empty text. Full response: {response}")
                    raise ValueError("Gemini generated an empty script.")
                generated_scripts.append(generated_text.strip())
            
            current_app.logger.info(f"Successfully processed chunk {i+1}/{len(chunks)}")
            
        except Exception as e:
            current_app.logger.error(f"Error processing chunk {i+1}/{len(chunks)}: {e}")
            if "504" in str(e) or "Deadline Exceeded" in str(e):
                raise ValueError("The request took too long to process. Please try with a shorter text or split it into smaller parts.")
            raise ValueError(f"Failed to generate script for chunk {i+1}: {str(e)}")

    # Combine all generated scripts with overlap handling
    final_script = []
    for i, script in enumerate(generated_scripts):
        if i > 0:
            # Find the last complete sentence in the previous script
            prev_sentences = script.split('.')
            if len(prev_sentences) > 1:
                # Remove the last sentence if it might be incomplete due to overlap
                final_script.append('.'.join(prev_sentences[:-1]) + '.')
        else:
            final_script.append(script)

    # Add the last chunk completely
    if generated_scripts:
        final_script.append(generated_scripts[-1])

    combined_script = "\n\n".join(final_script)
    current_app.logger.info("Successfully combined all chunks into final script")
    return combined_script


@presentation_bp.route('/convert_presentation_to_script', methods=['GET', 'POST'])
def convert_presentation_to_script():
    if request.method == 'POST':
        presentation_file = request.files.get('presentation_file')
        script_style = request.form.get('script_style')
        output_language = request.form.get('output_language', 'en')
        
        if not presentation_file or presentation_file.filename == '':
            return render_template('convert_presentation.html', error="No presentation file selected.")
        
        if not script_style:
            return render_template('convert_presentation.html', error="No script style selected.")

        allowed_extensions = {'pptx', 'pdf', 'md', 'txt'}
        file_ext = presentation_file.filename.rsplit('.', 1)[1].lower() if '.' in presentation_file.filename else ''
        
        if file_ext not in allowed_extensions:
            return render_template('convert_presentation.html', error="Invalid file type. Please upload a .pptx, .pdf, .md, or .txt file.")

        try:
            current_app.logger.info(f"Processing presentation file: {presentation_file.filename}, style: {script_style}, output_language: {output_language}")
            file_stream = io.BytesIO(presentation_file.read())
            
            if file_ext == 'pptx':
                presentation_text = extract_text_from_pptx(file_stream)
            elif file_ext == 'pdf':
                presentation_text = extract_text_from_pdf(file_stream)
            elif file_ext in {'md', 'txt'}:
                file_stream.seek(0)
                presentation_text = file_stream.read().decode('utf-8')
                current_app.logger.info(f"Extracted text from {file_ext}: {presentation_text}")
            else:
                presentation_text = ''
            
            if not presentation_text.strip():
                return render_template('convert_presentation.html', error="Could not extract any text from the presentation, or the presentation is empty.")

            current_app.logger.info(f"Extracted text length: {len(presentation_text)}")
            generated_script = generate_script_with_gemini(presentation_text, script_style, output_language)
            
            return render_template('convert_presentation.html', generated_script=generated_script, output_language=output_language)

        except ValueError as ve:
            current_app.logger.error(f"ValueError during presentation conversion: {ve}")
            return render_template('convert_presentation.html', error=str(ve))
        except Exception as e:
            current_app.logger.error(f"Unexpected error during presentation conversion: {e}", exc_info=True)
            return render_template('convert_presentation.html', error=f"An unexpected error occurred: {str(e)}")

    return render_template('convert_presentation.html')